from pptx import Presentation
from pptx.util import Pt
import json

prs = Presentation(r"C:\Users\munoz\sensai_agents\outputs\liver_timeline\HCC-TX-Research-Timeline-v5.pptx")

for slide_num, slide in enumerate(prs.slides, 1):
    print(f"\n=== SLIDE {slide_num} ===")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f"  [{shape.shape_id}] {shape.name}: {text[:120]}")
