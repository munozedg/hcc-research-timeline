from pptx import Presentation
import sys

prs = Presentation(r"C:\Users\munoz\sensai_agents\outputs\liver_timeline\HCC-TX-Research-Timeline-v5.pptx")

output_lines = []
output_lines.append(f"Total slides: {len(prs.slides)}")

for slide_num, slide in enumerate(prs.slides, 1):
    output_lines.append(f"\n=== SLIDE {slide_num} ===")
    output_lines.append(f"  Total shapes: {len(slide.shapes)}")
    for shape in slide.shapes:
        output_lines.append(f"  --- Shape [{shape.shape_id}] '{shape.name}' (has_text={shape.has_text_frame}) ---")
        if shape.has_text_frame:
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if text:
                    output_lines.append(f"    para[{para_idx}]: {text}")

result = "\n".join(output_lines)
# Write to file
with open(r"C:\Users\munoz\sensai_agents\outputs\liver_timeline\pptx_inspection.txt", "w", encoding="utf-8") as f:
    f.write(result)

print("Done. Written to pptx_inspection.txt")
print(f"Total lines: {len(output_lines)}")
