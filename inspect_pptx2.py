from pptx import Presentation

prs = Presentation(r"C:\Users\munoz\sensai_agents\outputs\liver_timeline\HCC-TX-Research-Timeline-v5.pptx")

print(f"Total slides: {len(prs.slides)}")

for slide_num, slide in enumerate(prs.slides, 1):
    print(f"\n=== SLIDE {slide_num} ===")
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(f"  --- Shape [{shape.shape_id}] '{shape.name}' ---")
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if text:
                    print(f"    para[{para_idx}]: {text}")
